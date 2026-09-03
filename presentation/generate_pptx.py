import os
import sys
import re
from html import unescape
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Add current directory to path
current_dir = os.path.dirname(os.path.abspath(__file__))
if current_dir not in sys.path:
    sys.path.insert(0, current_dir)

import slides_sec0, slides_sec1, slides_sec2, slides_sec3, slides_sec4, slides_sec5, slides_sec6

# 16:9 Widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color Palette
PEA_PURPLE = RGBColor(107, 63, 160)       # #6B3FA0
PEA_DARK = RGBColor(74, 37, 116)          # #4A2574
PEA_DEEP = RGBColor(47, 19, 77)           # #2F134D
PEA_LIGHT_BG = RGBColor(248, 246, 252)    # #F8F6FC
STATUS_RED = RGBColor(166, 55, 58)        # #A6373A
STATUS_RED_BG = RGBColor(253, 235, 235)   # #FDEBEB
STATUS_GREEN = RGBColor(39, 107, 71)      # #276B47
STATUS_GREEN_BG = RGBColor(234, 245, 238) # #EAF5EE
STATUS_AMBER = RGBColor(184, 99, 30)      # #B8631E
STATUS_AMBER_BG = RGBColor(254, 244, 235) # #FEF4EB
STATUS_BLUE = RGBColor(29, 93, 166)       # #1D5DA6
STATUS_BLUE_BG = RGBColor(234, 242, 251)  # #EAF2FB
TEXT_MAIN = RGBColor(33, 37, 41)          # #212529
TEXT_MUTED = RGBColor(108, 117, 125)      # #6C757D
BORDER_LIGHT = RGBColor(220, 215, 230)    # #DCD7E6
WHITE = RGBColor(255, 255, 255)
DARK_CARD = RGBColor(24, 20, 36)          # #181424

FONT_HEADING = "Kanit"
FONT_BODY = "IBM Plex Sans Thai"
FONT_MONO = "Consolas"

SECTION_NAMES = {
    0: "0. ภาพรวมโครงการ",
    1: "1. ปัญหาและสาเหตุรากเหง้า (Root Causes)",
    2: "2. แนวคิดและหลักการ Agentic AI",
    3: "3. สถาปัตยกรรมระบบ & เทคโนโลยี",
    4: "4. ช่องทางบริการ & ระบบปลั๊กอิน",
    5: "5. ตัวอย่างการทำงานจริง (Use Cases)",
    6: "6. ผลลัพธ์ การประเมิน & แผนงาน"
}

def clean_html_text(raw_html):
    # Remove script and style tags
    clean = re.sub(r'<(script|style).*?</\1>', '', raw_html, flags=re.DOTALL)
    # Replace breaks and block tags with newlines
    clean = re.sub(r'<(br|/p|/div|/h1|/h2|/h3|/h4|/li|/tr)>', '\n', clean)
    # Remove all remaining HTML tags
    clean = re.sub(r'<[^>]+>', '', clean)
    # Unescape HTML entities
    clean = unescape(clean)
    # Normalize whitespace while preserving linebreaks
    lines = [line.strip() for line in clean.splitlines() if line.strip()]
    return lines

def extract_cards(html_content):
    cards = []
    # Match card containers or major divs
    pattern = r'<div class="[^"]*(?:card|box|item|layer|step|spec|stat|col|phase)[^"]*"[^>]*>(.*?)</div>'
    matches = re.findall(pattern, html_content, flags=re.DOTALL)
    for m in matches:
        lines = clean_html_text(m)
        if len(lines) >= 1:
            title = lines[0]
            body = lines[1:]
            cards.append((title, body))
    return cards

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_topbar(slide, slide_num, total_slides, section_id, rubric_text):
    # Topbar background
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.58))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.color.rgb = BORDER_LIGHT
    bar.line.width = Pt(1)

    # PEA Badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(0.11), Inches(1.8), Inches(0.36))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PEA_PURPLE
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "⚡ PEA ONE AGENT"
    p.font.name = FONT_HEADING
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(2.25), Inches(0.11), Inches(2.2), Inches(0.36))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "หัวข้อ 3: Agentic AI"
    p.font.name = FONT_BODY
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_MUTED

    # Section Indicator
    sec_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(0.11), Inches(3.4), Inches(0.36))
    sec_box.fill.solid()
    sec_box.fill.fore_color.rgb = PEA_LIGHT_BG
    sec_box.line.color.rgb = BORDER_LIGHT
    sec_box.line.width = Pt(0.5)
    tf = sec_box.text_frame
    p = tf.paragraphs[0]
    sec_name = SECTION_NAMES.get(section_id, f"หมวด {section_id}")
    p.text = f"📍 {sec_name}"
    p.font.name = FONT_BODY
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = PEA_DARK
    p.alignment = PP_ALIGN.CENTER

    # Rubric Badge
    rub_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(0.11), Inches(2.7), Inches(0.36))
    rub_box.fill.solid()
    if "เกณฑ์ C" in rubric_text:
        rub_box.fill.fore_color.rgb = STATUS_GREEN_BG
        text_col = STATUS_GREEN
    elif "เกณฑ์ B" in rubric_text:
        rub_box.fill.fore_color.rgb = PEA_LIGHT_BG
        text_col = PEA_PURPLE
    elif "เกณฑ์ A" in rubric_text:
        rub_box.fill.fore_color.rgb = STATUS_AMBER_BG
        text_col = STATUS_AMBER
    else:
        rub_box.fill.fore_color.rgb = PEA_LIGHT_BG
        text_col = PEA_DARK

    rub_box.line.color.rgb = BORDER_LIGHT
    rub_box.line.width = Pt(0.5)
    tf = rub_box.text_frame
    p = tf.paragraphs[0]
    p.text = rubric_text
    p.font.name = FONT_BODY
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = text_col
    p.alignment = PP_ALIGN.CENTER

    # Slide counter
    cnt_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.8), Inches(0.11), Inches(1.2), Inches(0.36))
    cnt_box.fill.solid()
    cnt_box.fill.fore_color.rgb = RGBColor(240, 240, 245)
    cnt_box.line.color.rgb = BORDER_LIGHT
    cnt_box.line.width = Pt(0.5)
    tf = cnt_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_num:02d} / {total_slides:02d}"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

def add_titles(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12.13), Inches(0.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_HEADING
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = PEA_DARK

    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.13), Inches(0.35))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.name = FONT_BODY
    p.font.size = Pt(11.5)
    p.font.color.rgb = TEXT_MUTED

def build_cover_slide(prs, s_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PEA_DEEP)

    tag_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.33), Inches(0.5))
    p = tag_box.text_frame.paragraphs[0]
    p.text = "⚡ PEA INNOVATION AWARDS 2026 — TOPIC 3: AGENTIC AI FOR CUSTOMER SERVICE"
    p.font.name = FONT_HEADING
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 166, 35) # Gold
    p.alignment = PP_ALIGN.CENTER

    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.33), Inches(1.4))
    p = t_box.text_frame.paragraphs[0]
    p.text = "PEA One Agent"
    p.font.name = FONT_HEADING
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(11.33), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = 'จาก "ตอบคำถาม" สู่ "ทำงานแทนลูกค้า" — Agentic AI สำหรับศูนย์บริการข้อมูลผู้ใช้ไฟฟ้า 1129'
    p.font.name = FONT_BODY
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(225, 215, 245)
    p.alignment = PP_ALIGN.CENTER

    cov_feats = [
        ("🧠 สมองเดียว 3 ช่องทาง", "Web Chat, Voice (Gemini Live) และ LINE ผ่าน Main Agent เดี่ยว"),
        ("🛡️ ปลอดภัยระดับโค้ด", "State Machine บังคับ Human-in-the-Loop ยืนยันก่อนเขียนข้อมูล ปราศจาก Hallucination"),
        ("🔌 Plugin Architecture", "ขยายระบบใหม่ผ่าน Declarative Manifest โดยไม่ต้องแก้โค้ดระบบหลัก"),
        ("💰 ประหยัด 16.4 ล้าน/ปี", "ลดต้นทุนจาก 9.2 เหลือ 3.0 บาท/สาย (~67.4%) รองรับสายพร้อมกันไม่จำกัด")
    ]
    for idx, (ftitle, fdesc) in enumerate(cov_feats):
        c_left = Inches(0.8 + idx * 2.95)
        c_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(4.3), Inches(2.8), Inches(2.2))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = RGBColor(60, 28, 95)
        c_box.line.color.rgb = RGBColor(120, 80, 170)
        tf = c_box.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        p1 = tf.paragraphs[0]
        p1.text = ftitle
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(245, 166, 35)
        p2 = tf.add_paragraph()
        p2.text = fdesc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = RGBColor(220, 215, 235)

    p_cred = slide.shapes.add_textbox(Inches(1.0), Inches(6.8), Inches(11.33), Inches(0.4)).text_frame.paragraphs[0]
    p_cred.text = "ทีมพัฒนา PEA One Agent | กุมภาพันธ์ 2569 | เวลาบรรยาย 10 นาที"
    p_cred.font.name = FONT_BODY
    p_cred.font.size = Pt(10.5)
    p_cred.font.color.rgb = RGBColor(180, 170, 205)
    p_cred.alignment = PP_ALIGN.CENTER
    return slide

def build_closing_slide(prs, s_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PEA_DEEP)

    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.33), Inches(1.2))
    p = t_box.text_frame.paragraphs[0]
    p.text = "PEA One Agent"
    p.font.name = FONT_HEADING
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.33), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "พร้อมแล้วสำหรับการยกระดับงานบริการลูกค้า กฟภ. สู่มาตรฐาน Agentic AI ระดับโลก"
    p.font.name = FONT_BODY
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 166, 35)
    p.alignment = PP_ALIGN.CENTER

    summary_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.6), Inches(10.33), Inches(2.6))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(60, 28, 95)
    summary_box.line.color.rgb = RGBColor(120, 80, 170)
    tf = summary_box.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    p1 = tf.paragraphs[0]
    p1.text = "🎯 3 เสาหลักความสำเร็จที่เราพิสูจน์แล้วในการแข่งขันนี้:"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

    pillars = [
        "1. ถูกต้องและปลอดภัยระดับโค้ด: Human-in-the-Loop State Machine + Full-Doc Grounding ปราศจาก Hallucination",
        "2. ขยายระบบไร้ขีดจำกัด: Declarative Plugin Architecture เพิ่มระบบใหม่ใน 6 ขั้นตอนโดยไม่แตะสมองหลัก",
        "3. ผลลัพธ์ทางธุรกิจจับต้องได้: ประหยัด ~16.4 ล้านบาท/ปี ยกระดับ SLA รับสายใน 10 วินาทีเป็น 95%+"
    ]
    for pil in pillars:
        p = tf.add_paragraph()
        p.text = pil
        p.font.name = FONT_BODY
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(225, 215, 245)

    p_end = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(11.33), Inches(0.5)).text_frame.paragraphs[0]
    p_end.text = "ขอขอบพระคุณคณะกรรมการทุกท่าน — เปิดรับข้อซักถามและคำแนะนำ (Q&A)"
    p_end.font.name = FONT_HEADING
    p_end.font.size = Pt(15)
    p_end.font.bold = True
    p_end.font.color.rgb = WHITE
    p_end.alignment = PP_ALIGN.CENTER
    return slide

def build_slide_13_topology(prs, s_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PEA_LIGHT_BG)
    add_topbar(slide, 13, 43, 3, s_data["rubric"])
    add_titles(slide, s_data["title"], s_data["subtitle"])

    flow_img_path = os.path.join(current_dir, "assets", "flow.jpg")
    if not os.path.exists(flow_img_path):
        flow_img_path = os.path.join(os.path.dirname(current_dir), "assets", "flow.jpg")

    if os.path.exists(flow_img_path):
        f_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.35))
        f_card.fill.solid()
        f_card.fill.fore_color.rgb = WHITE
        f_card.line.color.rgb = BORDER_LIGHT
        slide.shapes.add_picture(flow_img_path, Inches(0.95), Inches(1.8), width=Inches(11.4), height=Inches(4.15))

    arch_callouts = [
        ("🌐 Multi-Channel ➔ FastAPI", "Web Chat, Voice WebSocket, LINE Webhook ใน Process เดียว ไม่แยกส่วน"),
        ("🧠 Main Agent (Gateway)", "สมองเดี่ยวคุม Bounded Loop ≤ 12 Steps ปราศจาก Sub-agent sprawl"),
        ("🔌 Plugin Runtime System", "ขยายระบบผ่าน Declarative plugin.yaml ตรวจสอบ Fail-Closed 7 ขั้นตอน"),
        ("🛡️ Safe Write State Machine", "Prepare ➔ Review ➔ มนุษย์กดยืนยัน ➔ Submit ป้องกันเสี่ยง 100%")
    ]
    for idx, (ctitle, cdesc) in enumerate(arch_callouts):
        c_left = Inches(0.8 + idx * 2.95)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(6.2), Inches(2.8), Inches(0.95))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = PEA_PURPLE if idx == 1 else BORDER_LIGHT
        card.line.width = Pt(1.5 if idx == 1 else 1)
        tf = card.text_frame
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.08)
        p1 = tf.paragraphs[0]
        p1.text = ctitle
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = PEA_DARK
        p2 = tf.add_paragraph()
        p2.text = cdesc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED
    return slide

def build_slide_26_line(prs, s_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PEA_LIGHT_BG)
    add_topbar(slide, 26, 43, 4, s_data["rubric"])
    add_titles(slide, s_data["title"], s_data["subtitle"])

    phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(4.8), Inches(5.2))
    phone.fill.solid()
    phone.fill.fore_color.rgb = WHITE
    phone.line.color.rgb = RGBColor(6, 199, 85)
    phone.line.width = Pt(2)
    tf = phone.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.text = "LINE OA: PEA Official Service"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(11.5)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(6, 199, 85)

    p2 = tf.add_paragraph()
    p2.text = "👤 ลูกค้า: แจ้งเรื่องร้องเรียนมิเตอร์ไฟฟ้าอ่านผิดครับ"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MAIN

    p3 = tf.add_paragraph()
    p3.text = "⚡ AI: PEA One Agent ได้ร่างเคส VOC ให้ท่านแล้ว กรุณากดยืนยันผ่านปุ่มด้านล่าง [Simulation]"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(9.5)
    p3.font.color.rgb = PEA_DARK

    p4 = tf.add_paragraph()
    p4.text = "🔘 [ยืนยันเปิดเคส VOC]    ❌ [ยกเลิก]"
    p4.font.name = FONT_HEADING
    p4.font.size = Pt(10.5)
    p4.font.bold = True
    p4.font.color.rgb = RGBColor(6, 199, 85)

    richmenu_path = os.path.join(current_dir, "assets", "line_richmenu.jpg")
    if not os.path.exists(richmenu_path):
        richmenu_path = os.path.join(os.path.dirname(current_dir), "assets", "line_richmenu.jpg")

    if os.path.exists(richmenu_path):
        slide.shapes.add_picture(richmenu_path, Inches(0.95), Inches(4.1), width=Inches(4.5), height=Inches(2.7))

    line_rules = [
        ("🔒 X-Line-Signature Verification", "ตรวจสอบลายเซ็น HMAC-SHA256 ของ Raw Request Body ทุกครั้ง หากไม่ถูกต้องจะตัดทิ้งด้วย HTTP 403 Forbidden ทันทีแบบ Fail-Closed"),
        ("🚫 No Free-Text Confirmation", "การตีความคำว่า 'ใช่/ตกลง/เอาเลย' จากแชตถูกตัดออกเพื่อความปลอดภัย ต้องกดปุ่ม Postback ที่มี ID กำกับเท่านั้น ป้องกันโมเดลตีความผิดพลาด"),
        ("⚡ Fast Return & Background Loop", "ตอบกลับ HTTP 200 ให้ LINE ภายใน 1 วินาทีเพื่อป้องกัน Timeout แล้วประมวลผล Agent Loop ใน Background พร้อมแสดงสถานะ"),
        ("🏷️ Simulation Tagging", "แสดงป้าย [Simulation] ทุกครั้งเมื่อผลลัพธ์เชื่อมต่อกับ Backend จำลอง เพื่อความโปร่งใสต่อผู้ใช้และคณะกรรมการ")
    ]
    for idx, (rtitle, rdesc) in enumerate(line_rules):
        r_top = Inches(1.75 + idx * 1.32)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), r_top, Inches(6.6), Inches(1.22))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_LIGHT
        tf = card.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.12)
        p1 = tf.paragraphs[0]
        p1.text = rtitle
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(11.5)
        p1.font.bold = True
        p1.font.color.rgb = PEA_DARK
        p2 = tf.add_paragraph()
        p2.text = rdesc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_MUTED
    return slide

def build_slide_38_pilot(prs, s_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PEA_LIGHT_BG)
    add_topbar(slide, 38, 43, 6, s_data["rubric"])
    add_titles(slide, s_data["title"], s_data["subtitle"])

    # Unit 1 Card
    u1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(5.7), Inches(4.55))
    u1.fill.solid()
    u1.fill.fore_color.rgb = WHITE
    u1.line.color.rgb = PEA_PURPLE
    u1.line.width = Pt(2)
    tf1 = u1.text_frame
    tf1.margin_left = Inches(0.2)
    tf1.margin_top = Inches(0.15)
    p = tf1.paragraphs[0]
    p.text = "🏢 หน่วยงานนำร่องที่ 1 (Voice / 24 ชม.)"
    p.font.name = FONT_HEADING
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PEA_DARK

    u1_items = [
        "ชื่อหน่วยงาน: [ระบุ เช่น ศูนย์บริการข้อมูลผู้ใช้ไฟฟ้า 1129]",
        "ขอบเขตงาน: [ระบุ เช่น ทดสอบ Voice Agent รับแจ้งเหตุไฟฟ้าขัดข้องและสอบถามข้อมูล 12 เขต]",
        "--- ผลการวัดผลเชิงปริมาณ (Quantitative Metrics) ---",
        "• % Abandon Call: เดิม 10.29% ➔ ผลทดลอง [ใส่ตัวเลข เช่น < 3.5%]",
        "• Speed of Answer: เดิม 46 วินาที ➔ ผลทดลอง [ใส่ตัวเลข เช่น 8 วินาที]",
        "• ปริมาณสายทดสอบ: [ระบุ เช่น 1,500 สาย] (รองรับพร้อมกันไม่จำกัด)",
        "• Task Completion Rate: [ระบุ เช่น 93.8% ปิดเคสได้เบ็ดเสร็จ]",
        "--- เสียงสะท้อนเชิงคุณภาพ (User Feedback) ---",
        '"[พิมพ์ใส่ความคิดเห็นจริง เช่น: เจ้าหน้าที่หน้างานลดความตึงเครียดช่วงฝนตกฟ้าคะนอง ระบบช่วยคัดกรองสายแจ้งเหตุซ้ำได้อย่างแม่นยำ ทำให้เจ้าหน้าที่มุ่งแก้ไขจุดวิกฤตได้รวดเร็วขึ้น]"'
    ]
    for item in u1_items:
        p = tf1.add_paragraph()
        p.text = item
        p.font.name = FONT_BODY
        p.font.size = Pt(9.5)
        if "---" in item:
            p.font.bold = True
            p.font.color.rgb = PEA_PURPLE
        elif "ผลทดลอง" in item:
            p.font.bold = True
            p.font.color.rgb = STATUS_GREEN
        elif '"' in item:
            p.font.italic = True
            p.font.color.rgb = RGBColor(80, 70, 100)

    # Unit 2 Card
    u2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.75), Inches(5.7), Inches(4.55))
    u2.fill.solid()
    u2.fill.fore_color.rgb = WHITE
    u2.line.color.rgb = STATUS_BLUE
    u2.line.width = Pt(2)
    tf2 = u2.text_frame
    tf2.margin_left = Inches(0.2)
    tf2.margin_top = Inches(0.15)
    p = tf2.paragraphs[0]
    p.text = "🏢 หน่วยงานนำร่องที่ 2 (Web & LINE / หน้าร้าน)"
    p.font.name = FONT_HEADING
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = STATUS_BLUE

    u2_items = [
        "ชื่อหน่วยงาน: [ระบุ เช่น การไฟฟ้าส่วนภูมิภาค สาขาเมืองเชียงใหม่ / กฟฟ.สาขา...]",
        "ขอบเขตงาน: [ระบุ เช่น ทดสอบระบบ Web Chat และ LINE OA สำหรับรับเรื่องร้องเรียน VOC และขอใช้ไฟฟ้า]",
        "--- ผลการวัดผลเชิงปริมาณ (Quantitative Metrics) ---",
        "• ระยะเวลากรอกคำร้อง: เดิม 15-20 นาที ➔ ผลทดลอง [ใส่ตัวเลข เช่น 45 วินาที]",
        "• ความถูกต้องข้อมูล (CA): เดิมต้องตรวจซ้ำ ➔ ผลทดลอง [ใส่ตัวเลข เช่น 100% ผ่าน]",
        "• ปริมาณเคสทดสอบ: [ระบุ เช่น 450 คำร้อง]",
        "• ความพึงพอใจลูกค้า (CSAT): [ระบุ เช่น 4.85 / 5.0 คะแนน]",
        "--- เสียงสะท้อนเชิงคุณภาพ (User Feedback) ---",
        '"[พิมพ์ใส่ความคิดเห็นจริง เช่น: ผู้ใช้ไฟประทับใจที่สามารถแนบพิกัดและรูปถ่ายผ่าน LINE ได้ทันที ระบบสร้าง Pending Action ให้ตรวจทานก่อนส่งเข้า VOC อัตโนมัติ ประหยัดเวลามาก]"'
    ]
    for item in u2_items:
        p = tf2.add_paragraph()
        p.text = item
        p.font.name = FONT_BODY
        p.font.size = Pt(9.5)
        if "---" in item:
            p.font.bold = True
            p.font.color.rgb = STATUS_BLUE
        elif "ผลทดลอง" in item:
            p.font.bold = True
            p.font.color.rgb = STATUS_GREEN
        elif '"' in item:
            p.font.italic = True
            p.font.color.rgb = RGBColor(60, 80, 100)

    # Bottom Scale-out strip
    b_strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.65))
    b_strip.fill.solid()
    b_strip.fill.fore_color.rgb = WHITE
    b_strip.line.color.rgb = BORDER_LIGHT
    tf = b_strip.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "🚀 แผนขยายผลสู่ 12 เขต: Shadow Mode (2 สัปดาห์) ➔ Limited Live (4 สัปดาห์) ➔ Nationwide Rollout ครอบคลุม 21 ล้านราย | 📱 Demo: http://127.0.0.1:8000"
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = PEA_PURPLE
    return slide

def build_generic_slide(prs, s_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PEA_LIGHT_BG)
    add_topbar(slide, s_data["id"], 43, s_data["section"], s_data["rubric"])
    add_titles(slide, s_data["title"], s_data["subtitle"])

    content = s_data["content"]
    text_lines = clean_html_text(content)

    # If slide contains code block or pre
    has_code = "<pre>" in content or "<code>" in content or "class=\"code-body\"" in content
    
    # Try extracting structured cards
    cards = extract_cards(content)
    
    # If we found 2 to 6 cards, render them cleanly in columns
    if len(cards) >= 2 and len(cards) <= 6:
        n = len(cards)
        if n == 2:
            col_w = Inches(5.7)
            for idx, (ctitle, clines) in enumerate(cards):
                left = Inches(0.8 + idx * 6.0)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.75), col_w, Inches(5.2))
                card.fill.solid()
                card.fill.fore_color.rgb = WHITE
                card.line.color.rgb = BORDER_LIGHT
                tf = card.text_frame
                tf.margin_left = Inches(0.25)
                tf.margin_top = Inches(0.2)
                p = tf.paragraphs[0]
                p.text = ctitle
                p.font.name = FONT_HEADING
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = PEA_DARK
                for line in clines[:12]:
                    p2 = tf.add_paragraph()
                    p2.text = f"• {line}" if not line.startswith("•") else line
                    p2.font.name = FONT_BODY
                    p2.font.size = Pt(10.5)
                    p2.font.color.rgb = TEXT_MAIN
        elif n == 3:
            col_w = Inches(3.75)
            for idx, (ctitle, clines) in enumerate(cards):
                left = Inches(0.8 + idx * 4.0)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.75), col_w, Inches(5.2))
                card.fill.solid()
                card.fill.fore_color.rgb = WHITE
                card.line.color.rgb = BORDER_LIGHT
                tf = card.text_frame
                tf.margin_left = Inches(0.2)
                tf.margin_top = Inches(0.2)
                p = tf.paragraphs[0]
                p.text = ctitle
                p.font.name = FONT_HEADING
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = PEA_DARK
                for line in clines[:10]:
                    p2 = tf.add_paragraph()
                    p2.text = f"• {line}" if not line.startswith("•") else line
                    p2.font.name = FONT_BODY
                    p2.font.size = Pt(10)
                    p2.font.color.rgb = TEXT_MAIN
        else: # 4 to 6 cards in 2 rows
            for idx, (ctitle, clines) in enumerate(cards[:6]):
                row = idx // 3
                col = idx % 3
                left = Inches(0.8 + col * 4.0)
                top = Inches(1.75 + row * 2.65)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.8), Inches(2.45))
                card.fill.solid()
                card.fill.fore_color.rgb = WHITE
                card.line.color.rgb = BORDER_LIGHT
                tf = card.text_frame
                tf.margin_left = Inches(0.2)
                tf.margin_top = Inches(0.15)
                p = tf.paragraphs[0]
                p.text = ctitle
                p.font.name = FONT_HEADING
                p.font.size = Pt(11.5)
                p.font.bold = True
                p.font.color.rgb = PEA_DARK
                for line in clines[:5]:
                    p2 = tf.add_paragraph()
                    p2.text = f"• {line}" if not line.startswith("•") else line
                    p2.font.name = FONT_BODY
                    p2.font.size = Pt(9.5)
                    p2.font.color.rgb = TEXT_MAIN
    else:
        # 2-column generic layout
        mid = len(text_lines) // 2
        col1_lines = text_lines[:mid] if mid > 0 else text_lines
        col2_lines = text_lines[mid:] if mid > 0 else []

        card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(5.7), Inches(5.2))
        card1.fill.solid()
        card1.fill.fore_color.rgb = WHITE
        card1.line.color.rgb = BORDER_LIGHT
        tf1 = card1.text_frame
        tf1.margin_left = Inches(0.25)
        tf1.margin_top = Inches(0.2)
        if col1_lines:
            p = tf1.paragraphs[0]
            p.text = col1_lines[0]
            p.font.name = FONT_HEADING
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = PEA_DARK
            for line in col1_lines[1:14]:
                p2 = tf1.add_paragraph()
                p2.text = f"• {line}" if not line.startswith("•") else line
                p2.font.name = FONT_BODY
                p2.font.size = Pt(10.5)
                p2.font.color.rgb = TEXT_MAIN

        if col2_lines:
            card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.75), Inches(5.7), Inches(5.2))
            card2.fill.solid()
            card2.fill.fore_color.rgb = WHITE
            card2.line.color.rgb = BORDER_LIGHT
            tf2 = card2.text_frame
            tf2.margin_left = Inches(0.25)
            tf2.margin_top = Inches(0.2)
            p = tf2.paragraphs[0]
            p.text = col2_lines[0]
            p.font.name = FONT_HEADING
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = PEA_DARK
            for line in col2_lines[1:14]:
                p2 = tf2.add_paragraph()
                p2.text = f"• {line}" if not line.startswith("•") else line
                p2.font.name = FONT_BODY
                p2.font.size = Pt(10.5)
                p2.font.color.rgb = TEXT_MAIN

    return slide

def build_all_presentation():
    all_slides = (
        slides_sec0.get_slides() +
        slides_sec1.get_slides() +
        slides_sec2.get_slides() +
        slides_sec3.get_slides() +
        slides_sec4.get_slides() +
        slides_sec5.get_slides() +
        slides_sec6.get_slides()
    )

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print(f"Generating full 43 slides PowerPoint presentation...")

    for idx, s_data in enumerate(all_slides):
        s_id = s_data["id"]
        print(f"Rendering slide {s_id}/43: {s_data['title'][:35]}...")
        if s_id == 1:
            build_cover_slide(prs, s_data)
        elif s_id == 13:
            build_slide_13_topology(prs, s_data)
        elif s_id == 26:
            build_slide_26_line(prs, s_data)
        elif s_id == 38:
            build_slide_38_pilot(prs, s_data)
        elif s_id == 43:
            build_closing_slide(prs, s_data)
        else:
            build_generic_slide(prs, s_data)

    output_path = os.path.join(current_dir, "pea_one_agent_presentation.pptx")
    prs.save(output_path)
    print(f"🎉 Successfully built all 43 slides in {output_path} ({os.path.getsize(output_path):,} bytes)")

if __name__ == "__main__":
    build_all_presentation()
